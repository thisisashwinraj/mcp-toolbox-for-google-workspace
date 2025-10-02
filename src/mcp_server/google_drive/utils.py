# Copyright 2025 Ashwin Raj
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import io
import json
import logging
import functools

import xlrd
import openpyxl
from docx import Document
from PyPDF2 import PdfReader
from bs4 import BeautifulSoup
from pptx import Presentation

from googleapiclient.errors import HttpError

logger = logging.getLogger(__name__)


def parse_file_content(mime_type: str, content_bytes: bytes) -> str:
    def safe_decode(b: bytes) -> str:
        return b.decode("utf-8-sig", errors="replace")

    text_content = ""

    if mime_type == "text/html":
        text_content = BeautifulSoup(
            safe_decode(content_bytes), "html.parser"
        ).get_text()

    elif mime_type == "application/pdf":
        pdf = PdfReader(io.BytesIO(content_bytes))

        text_content = "\n".join(
            [page.extract_text() or "" for page in pdf.pages]
        )

    elif mime_type == "application/json":
        try:
            obj = json.loads(safe_decode(content_bytes))
            text_content = json.dumps(obj, indent=2)

        except Exception:
            text_content = safe_decode(content_bytes)

    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(io.BytesIO(content_bytes))
        text_content = "\n".join([p.text for p in doc.paragraphs])

    elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        wb = openpyxl.load_workbook(io.BytesIO(content_bytes), data_only=True)

        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text_content += "\t".join(
                    [str(cell) if cell else "" for cell in row]
                ) + "\n"

    elif mime_type == "application/vnd.ms-excel":
        book = xlrd.open_workbook(file_contents=content_bytes)

        for sheet in book.sheets():
            for row_idx in range(sheet.nrows):
                row = sheet.row_values(row_idx)

                text_content += "\t".join([str(cell) for cell in row]) + "\n"

    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        prs = Presentation(io.BytesIO(content_bytes))

        for i, slide in enumerate(prs.slides):
            text_content += f"\n--- Slide {i+1} ---\n"

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content += shape.text + "\n"

    else:
        text_content = safe_decode(content_bytes)

    return text_content

def handle_google_drive_exceptions(func):
    @functools.wraps(func)
    async def wrapper(*args, **kwargs):
        try:
            return await func(*args, **kwargs)

        except HttpError as http_error:
            status = http_error.resp.status
            reason = http_error._get_reason()
            
            if status == 400:
                message = "Invalid request or file metadata might be incorrect"
            
            elif status == 401:
                message = (
                    "Unauthorized access. "
                    "Check if the credentials are valid or expired."
                )

            elif status == 403:
                message = (
                    "Permission denied. "
                    "Ensure your OAuth scope includes drive access and that "
                    "authenticated user has permission to perform this action."
                )
            
            elif status == 404:
                message = (
                    "Resource not found. "
                    "The file, folder, or drive may not exist or was deleted."
                )

            elif status == 409:
                message = (
                    "Conflict error. "
                    "This could be due to duplicate operations."
                )

            elif status == 410:
                message ="The resource is no longer available."

            elif status == 412:
                message = (
                    "Precondition failed. "
                    "Try syncing again or verify versioning headers."
                )

            elif status == 429:
                message = (
                    "Quota exceeded. Too many requests. "
                    "Try again later or use exponential backoff."
                )
            
            elif status in [500, 503]:
                message = (
                    "The Google Drive service is temporarily unavailable. "
                    "Please retry after some time."
                )
            
            else:
                message = f"Unexpected error with the Drive API: {reason}"

            logger.error(f"[Status Code {status}]: {message}", exc_info=True)

            return {
                "status": "error", 
                "message": message
            }

        except Exception as error:
            logger.error("An unexpected error occured", exc_info=True)

            return {
                "status": "error", 
                "message": f"An unexpected error occurred: {error}", 
            }
    
    return wrapper
